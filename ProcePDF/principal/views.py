from django.shortcuts import render
from PIL import Image
import pytesseract
from django.http import JsonResponse
from django.core.files.storage import FileSystemStorage
from django.conf import settings
import os
import fitz  
from nltk.corpus import stopwords
from nltk.tokenize import word_tokenize, sent_tokenize
from nltk.probability import FreqDist
import nltk

from wordcloud import WordCloud
from django.http import HttpResponse
import matplotlib.pyplot as plt
from io import BytesIO
import base64
import requests
import zipfile

from django.shortcuts import get_object_or_404
from django.http import FileResponse


nltk.download('punkt')
nltk.download('stopwords')

#------------------------------------------------------------- FUNCIONES PARA PASAR A OTRAS VENTANAS -----------------------------------------------------------
def index(request):
    return render(request, 'principal/index.html')

def ventanaIdeas(request):
    #extracted_text = request.session.get('extracted_text', '')  
    #return render(request, 'principal/ventanaIdeas.html', {'extracted_text': extracted_text})
    return render(request, 'principal/ventanaIdeas.html')

def ventanaCollage(request):
    return render(request, 'principal/ventanaCollage.html')

def ventanaResumen(request):
    return render(request, 'principal/ventanaResumen.html')

def ventanaImagenes(request):
    return render(request, 'principal/ventanaImagenes.html')

def ventanaPpt(request):
    return render(request, 'principal/ventanaPpt.html')


#----------------------------------------------------------------- FUNCIONES DE OPERACIONES -----------------------------------------------------------------
#Extrae el texto del archivo PDF
def extraerTexto(request):
    if request.method == 'POST' and request.FILES['pdf_file']:
        pdf_file = request.FILES['pdf_file']
        fs = FileSystemStorage()
        filename = fs.save(pdf_file.name, pdf_file)
        pdf_path = fs.url(filename)
        text = extraerTextoImagen(pdf_path)       
        request.session['extracted_text'] = text  # Guarda el texto en la sesión para que esté disponible en la vista de 'generarResumen'

        # Eliminar el archivo PDF después de obtener el texto
        full_path = os.path.join(settings.MEDIA_ROOT, pdf_path.strip('/'))
        if os.path.exists(full_path):
            os.remove(full_path)
        return JsonResponse({'text': text})
    return JsonResponse({'error': 'No se proporcionó un archivo PDF válido'})


#Extrae el texto de las imagenes del PDF
def extraerTextoImagen(pdf_path):
    image_dir = os.path.join(settings.MEDIA_ROOT, 'pdf_images')
    os.makedirs(image_dir, exist_ok=True)
    text = ''
    full_path = os.path.join(settings.MEDIA_ROOT, pdf_path.strip('/'))

    if os.path.exists(full_path):
        pdf = fitz.open(full_path)
        for page_num in range(pdf.page_count):
            page = pdf[page_num]
            image_list = page.get_pixmap()
            image_path = os.path.join(image_dir, f'page_{page_num + 1}.png')
            image_list.save(image_path)

            # Realiza la extracción de texto de la imagen con Pytesseract
            #extracted_text = pytesseract.image_to_string(Image.open(image_path))
            extracted_text = pytesseract.image_to_string(Image.open(image_path), lang='spa')

            text += extracted_text
        #eliminarImagen()
        return text
    else:
        return 'Archivo PDF no encontrado'


# Elimina todos los archivos en el directorio de imagenes
def eliminarImagen():
    image_dir = os.path.join(settings.MEDIA_ROOT, 'pdf_images')
    if os.path.exists(image_dir):
        for file in os.listdir(image_dir):
            file_path = os.path.join(image_dir, file)
            try:
                if os.path.isfile(file_path):
                    os.unlink(file_path)
            except Exception as e:
                print(f"Error al eliminar {file_path}: {e}")
        os.rmdir(image_dir)


#Genera un resumen del texto extraido
def generarResumen(request):
    idioma = set(stopwords.words("spanish")) #Idioma del texto
    text = request.session.get('extracted_text', '') #Texto de la ventana inicial
    palabras = word_tokenize(text)

    # Se crea un diccionario para crear una tabla de frecuencias de las palabras.
    frecuenciaTabla = dict()
    for palabra in palabras:
        palabra = palabra.lower() #formato de palabras
        if palabra in idioma: 
            continue 
        if palabra in frecuenciaTabla: # Si la palabra ya se encuentra en la tabla frecuencia se le Suma 1 a la posición donde se encuentra la palabra.
            frecuenciaTabla[palabra] += 1 
        else:
            frecuenciaTabla[palabra] = 1 # Sino, la palabra en la TF va a ser igual a 1.
    
    # Crea una variable y un diccionario.
    oraciones = sent_tokenize(text)
    oracionValor = dict()

    #Se resorren las oraciones que se encuentran en el texto.
    for oracion in oraciones:
        for palabra, freq in frecuenciaTabla.items():
            if palabra in oracion.lower(): 
                if oracion in oracionValor: 
                    oracionValor[oracion] += freq # Se le sume 1 al número de frecuencia en la posición de la oración del sV.
                else:
                    oracionValor[oracion] = freq # Sino, que el valor de la posición de la oración sea igual a la frecuencia.

    sumaValores = 0
    for oracion in oracionValor:
        sumaValores += oracionValor[oracion] # Se suma 1 al valor de la Oración en su respectiva posición       
    promedio = int(sumaValores/ len(oracionValor)) # Divide la suma de valores en la total de oraciones valorizadas. 

    #Si la oración está las oraciones Valorizadas y la posición de la oración es mayor que 1.2 veces el promedio
    resumen = ''
    for oracion in oraciones:
        if (oracion in oracionValor) and (oracionValor[oracion] > (1.2 * promedio)):    
            resumen += " " + oracion # El resumen va a agregar un espacio más la oración que aprobó la condición.
    return render(request, 'principal/ventanaResumen.html', {'text': text, 'resumen': resumen})

from django.utils.safestring import mark_safe


#Extrae las 5 ideas principales del texto
def generarIdeasPrincipales(request):
    if request.method == 'POST':
        text = request.session.get('extracted_text', '')  # Texto de la ventana inicial
        ideas = extraerIdeas(text, num_ideas=5)
        ideasF = ''
        for i, idea in enumerate(ideas, start=1):
            ideasF += str(i) +'- ' + idea + '<br>'
        ideasF = mark_safe(ideasF)  # Utiliza mark_safe para que se interpreten las etiquetas HTML
        return render(request, 'principal/ventanaIdeas.html', {'text': text, 'ideas_principales': ideasF})

def extraerIdeas(text, num_ideas=5):
    oraciones = sent_tokenize(text)     # Divide el texto en oraciones
    palabras = [word_tokenize(sentence.lower()) for sentence in oraciones]    # Tokeniza las palabras en cada oración

    # Calcula la frecuencia de las palabras
    todasPalabras = [word for sentence in palabras for word in sentence]
    fdist = FreqDist(todasPalabras)
    palabrasClave = [word for word, freq in fdist.most_common(num_ideas)]     # Encuentra las palabras clave más frecuentes
    # Encuentra las primeras 5 oraciones clave que contienen las palabras clave
    oracionesClave = []
    for sentence in oraciones:
        if any(keyword in sentence.lower() for keyword in palabrasClave):
            oracionesClave.append(sentence)
            if len(oracionesClave) >= num_ideas:
                break
    return oracionesClave


#Crea el collage de palabras
def generarCollage(request):
    if request.method == 'POST':
        texto = request.session.get('extracted_text', '')

        wc = WordCloud(width = 5000, height = 5000,
               stopwords = ["y", "a", "e", "i", "o", "u", "lo","ol","la", "al","el", "le","ella","ellos","ellas","él", "de","del", "una","un",
                            "en","es","eso","esa","esos","esas", "su","sus", "se", "por","era","ese","esa","mas","las","les","los","nol","´",
                            "'","para","pero","por","porque","por que","porqué","por qué","sin","yo","mi","tu","tú","usted","ustedes","casi",
                            "sean","si","no","talvez","quizas","quiso","muy","con","como","asi","así","ahi","ahí","alli","ay","hay","una",
                            "través","esto","este","esta","estos","estas","están","esten","tengan","tengo","tiene","te","ti","tí","deben","debe",
                            "debi","debí","hace","hacen","cuando","cuanto","cual","cuales","hice","echo","echos","antes","han","ha","he","echo",
                            "hacia","tanto","toda","todo","todas","todos","llega","que","mucho","eran","habían","había","habia","habian","fueron","fue",
                            "fui","hasta","mucho","mucha","muchos","muchas","pocos","pocas","dijo","dije","dij","junto","juntos","dejare","dejaré","nada",
                            "ñe","ser","soy","eres","ser","ya","tus","ir","vio","vi","visto","hummm", "ah", "aeeeh", "así que", "entonces","ok", "me",
                            "entiende", "más","mas", "sin embargo", "ver", "pues","puede","pueden","dio","da","fui","fue"
                            ])
        wc.generate(texto)

        # Cree una figura y muestre la imagen de la nube de palabras
        plt.figure(figsize=(30, 30))
        plt.imshow(wc, interpolation="bilinear")
        plt.axis("off")

        # Convierta la figura en una imagen PNG y luego en una cadena base64
        buffer = BytesIO()
        plt.savefig(buffer, format="png", bbox_inches='tight')
        image_data = buffer.getvalue()
        base64_image = base64.b64encode(image_data).decode("utf-8")      

        buffer.seek(0)
        image = Image.open(buffer)
        # Obtiene la ruta de la carpeta 'imagenes' en la raíz del proyecto
        images_dir = os.path.join(settings.BASE_DIR, 'imagenes')

        # Verifica si la carpeta 'imagenes' existe, y si no, créala
        if not os.path.exists(images_dir):
            os.makedirs(images_dir)
        # Guarda la imagen en la carpeta 'imagenes' con el nombre 'imageCollage.png'
        image_path = os.path.join(images_dir, 'imageCollage.png')
        image.save(image_path)

        return render(request, 'principal/ventanaCollage.html', {'collage_image': base64_image})


#Crea la busqueda de imagenes
def buscarImagenes(request): 
    #eliminarImagen()   
    access_key = "XXjQx3Qafdg7uqEiK96z3gcwc0PzDdeCw44NenAlAiw"  # Reemplaza con tu clave de acceso de Unsplash
    texto = request.session.get('extracted_text', '')
    search_term = palabraPrincipal(texto)
    url = f'https://api.unsplash.com/search/photos?query={search_term}&per_page=5'
    headers = {'Authorization': f'Client-ID {access_key}'}

    try:
        response = requests.get(url, headers=headers)
        response.raise_for_status()
        data = response.json()

        images = []
        for result in data['results']:
            images.append(result['urls']['regular'])

            image_url = result['urls']['regular']
            # Descarga la imagen
            response = requests.get(image_url, stream=True)
            response.raise_for_status()

            # Directorio donde se guardarán las imágenes
            image_directory = os.path.join(settings.MEDIA_ROOT, 'images')
            os.makedirs(image_directory, exist_ok=True)  # Crea el directorio si no existe

            # Guarda la imagen en tu servidor
            image_name = result['id'] + '.jpg'  # Puedes nombrar la imagen como quieras
            image_path = os.path.join(image_directory, image_name)

            with open(image_path, 'wb') as file:
                for chunk in response.iter_content(chunk_size=8192):
                    file.write(chunk)
        return JsonResponse({'images': images})
    except requests.exceptions.RequestException as e:
        return JsonResponse({'error': str(e)})

def palabraPrincipal(texto):# Obtener la palabra con mayor frecuencia
    wc = WordCloud(width=1400, height=1400,
                stopwords=[
                    "y", "a", "e", "i", "o", "u", "lo", "ol", "la", "al", "el", "le", "ella", "ellos", "ellas", "él", "de", "del","una", "un",
                    "en", "es", "eso", "esa", "esos", "esas", "su", "sus", "se", "por", "era", "ese", "esa", "mas", "las", "les", "los", "nol", "´",
                    "'", "para", "pero", "por", "porque", "por que", "porqué", "por qué", "sin", "yo", "mi", "tu", "tú", "usted", "ustedes", "casi",
                    "sean", "si", "no", "talvez", "quizas", "quiso", "muy", "con", "como", "asi", "así", "ahi", "ahí", "alli", "ay", "hay", "una",
                    "través", "esto", "este", "esta", "estos", "estas", "están", "esten", "tengan", "tengo", "tiene", "te", "ti", "tí", "deben", "debe",
                    "debi", "debí", "hace", "hacen", "cuando", "cuanto", "cual", "cuales", "hice", "echo", "echos", "antes", "han", "ha", "he", "echo",
                    "hacia", "tanto", "toda", "todo", "todas", "todos", "llega", "que", "mucho", "eran", "habían", "había", "habia", "habian", "fueron", "fue",
                    "fui", "hasta", "mucho", "mucha", "muchos", "muchas", "pocos", "pocas", "dijo", "dije", "dij", "junto", "juntos", "dejare", "dejaré", "nada",
                    "ñe", "ser", "soy", "eres", "ser", "ya", "tus", "ir", "vio", "vi", "visto", "hummm", "ah", "aeeeh", "así que", "entonces", "ok", "me",
                    "entiende", "más", "mas", "sin embargo", "ver", "pues", "puede", "pueden", "dio", "da", "fui", "fue"],max_words=1)
    wc.generate(texto)  
    palabraFrecuente = list(wc.words_.keys())[0]
    #print('TEXTO: ' + palabraFrecuente)
    return palabraFrecuente


#----------------------------------------------------------------- FUNCIONES DE DESCARGA -----------------------------------------------------------------
def descargarIdeas(request):
    ideas = request.POST.get('ideas')  # Obtén las ideas principales
    response = HttpResponse(content_type='text/plain')
    response['Content-Disposition'] = 'attachment; filename=ideas.txt'
    response.write(ideas)
    return response

def descargar_imagenes(request):
    images_directory = os.path.join(settings.MEDIA_ROOT, 'images')

    # Verifica si el directorio de imágenes existe
    if not os.path.exists(images_directory):
        return HttpResponse("No se encontraron imágenes para descargar.")

    # Crea un archivo ZIP en memoria
    buffer = BytesIO()
    with zipfile.ZipFile(buffer, 'w', zipfile.ZIP_DEFLATED) as zipf:
        # Recorre las imágenes en el directorio y agrégalas al archivo ZIP
        for root, dirs, files in os.walk(images_directory):
            for file in files:
                file_path = os.path.join(root, file)
                zipf.write(file_path, os.path.relpath(file_path, images_directory))

    # Configura la respuesta HTTP para el archivo ZIP
    response = HttpResponse(buffer.getvalue(), content_type='application/zip')
    response['Content-Disposition'] = 'attachment; filename=imagenes.zip'

    return response

def descargarResumen(request):
    resumen = request.POST.get('resumen')  # Obtén las ideas principales
    response = HttpResponse(content_type='text/plain')
    response['Content-Disposition'] = 'attachment; filename=resumen.txt'
    response.write(resumen)
    return response

def download_collage(request):
    images_directory = os.path.join(settings.MEDIA_ROOT, 'imagenes')
    image_filename = 'imageCollage.png'  # Nombre de la imagen a descargar

    # Verifica si el directorio de imágenes y la imagen existen
    if not os.path.exists(images_directory) or not os.path.exists(os.path.join(images_directory, image_filename)):
        return HttpResponse("No se encontró la imagen para descargar.")

    # Construye la ruta completa de la imagen
    image_path = os.path.join(images_directory, image_filename)

    # Abre el archivo de imagen y lo envía como respuesta sin cerrarlo
    image_file = open(image_path, 'rb')
    response = FileResponse(image_file, content_type='image/png')
    response['Content-Disposition'] = f'attachment; filename="{image_filename}"'

    return response




from django.http import FileResponse
from django.template.loader import render_to_string
from django.shortcuts import render
from django.core.files.storage import FileSystemStorage
from pptx import Presentation
from io import BytesIO
from django.templatetags.static import static


def generar_presentacion(request):
    # Obtiene el título de la primera diapositiva del usuario
    titulo = request.POST.get('titulo', 'Título de la presentación')

    # Crea una presentación PPT
    ppt = Presentation()

    # Crea la primera diapositiva con el título
    slide = ppt.slides.add_slide(ppt.slide_layouts[0])
    title = slide.shapes.title
    title.text = titulo

    # Recupera las ideas principales y las imágenes de las ideas
    ideas_principales = generarIdeasPrincipales(request)
    imagenes_ideas = buscarImagenesPPT(request)

    # Agrega una diapositiva por cada idea principal con su imagen
    for idea in ideas_principales:
        slide = ppt.slides.add_slide(ppt.slide_layouts[5])  # Puedes elegir un diseño de diapositiva adecuado
        title = slide.shapes.title
        title.text = "Idea Principal"

        # Define las coordenadas y dimensiones de la imagen
        left = 100  # Coordenada izquierda
        top = 100   # Coordenada superior
        width = 400  # Ancho
        height = 300  # Altura

        # Utiliza la ruta de la imagen por defecto desde los archivos estáticos de Django
        default_image_path = static('ruta_imagen_default.png')
        content = slide.shapes.add_picture(imagenes_ideas.get(idea, default_image_path), left, top, width, height)

        #ppppppcontent = slide.shapes.add_picture(imagenes_ideas.get(idea, 'ruta_imagen_default.png'), left, top, width, height)

        # Agrega la idea principal como texto adicional si es necesario
        tf = content.text_frame
        p = tf.add_paragraph()
        p.text = idea

    # Guarda la presentación en un archivo temporal
    ppt_file = BytesIO()
    ppt.save(ppt_file)

    # Envía el archivo como respuesta para que se pueda abrir en la ventana VentanaPpt
    ppt_file.seek(0)
    response = FileResponse(ppt_file, content_type='application/vnd.openxmlformats-officedocument.presentationml.presentation')
    response['Content-Disposition'] = f'attachment; filename=mi_presentacion.pptx'

    return response


#Crea la busqueda de imagenes
def buscarImagenesPPT(request):    
    access_key = "XXjQx3Qafdg7uqEiK96z3gcwc0PzDdeCw44NenAlAiw"  # Reemplaza con tu clave de acceso de Unsplash
    texto = request.session.get('extracted_text', '')
    search_term = palabraPrincipal(texto)
    url = f'https://api.unsplash.com/search/photos?query={search_term}&per_page=1'
    headers = {'Authorization': f'Client-ID {access_key}'}

    try:
        response = requests.get(url, headers=headers)
        response.raise_for_status()
        data = response.json()

        images = []
        for result in data['results']:
            images.append(result['urls']['regular'])
        # Registra las imágenes en la consola
        print("Imágenes encontradas:")
        for image_url in images:
            print(image_url)

        return JsonResponse({'images': images})
    except requests.exceptions.RequestException as e:
        # Registra errores en la consola
        print(f"Error al buscar imágenes: {str(e)}")
        return JsonResponse({'error': str(e)})
